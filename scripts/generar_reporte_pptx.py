import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
import plotly.graph_objects as go
from pptx import Presentation
from pptx.util import Inches
import os

# Crear carpeta para gráficos si no existe
os.makedirs("charts", exist_ok=True)

# Leer datos
df_summary = pd.read_csv("report_csvs/ResumenGlobal.csv")
df_sellers = pd.read_csv("report_csvs/Top10Sellers.csv")
df_dest = pd.read_csv("report_csvs/PorRegionDestino.csv")
df_origin = pd.read_csv("report_csvs/PorRegionOrigen.csv")
df_routes = pd.read_csv("report_csvs/Top10Rutas.csv")

# ------------------------------------------
# GENERAR GRÁFICOS COMO PNGs
# ------------------------------------------

# 1. Región Destino
plt.figure(figsize=(6,4))
sns.barplot(x='destination_region', y='under_rate_pct', data=df_dest, palette='pastel', edgecolor='gray')
plt.title("Under-Promise por Región Destino")
plt.ylabel("% Under-Promise")
plt.tight_layout()
plt.savefig("charts/destino.png")
plt.close()

# 2. Región Origen
plt.figure(figsize=(6,4))
sns.barplot(x='origin_region', y='under_rate_pct', data=df_origin, palette='pastel', edgecolor='gray')
plt.title("Under-Promise por Región Origen")
plt.ylabel("% Under-Promise")
plt.tight_layout()
plt.savefig("charts/origen.png")
plt.close()

# 3. Top Sellers
plt.figure(figsize=(6,4))
sns.barplot(y='seller_id', x='under_rate_pct', data=df_sellers.sort_values('under_rate_pct'), palette='Set2', edgecolor='gray')
plt.title("Top 10 Sellers por Under-Promise")
plt.xlabel("% Under-Promise")
plt.tight_layout()
plt.savefig("charts/sellers.png")
plt.close()

# 4. Heatmap Rutas
heatmap_data = df_routes.pivot(index='origin_region', columns='destination_region', values='under_rate_pct')
plt.figure(figsize=(6,4))
sns.heatmap(heatmap_data, annot=True, fmt=".1f", cmap="Blues")
plt.title("Heatmap Rutas Origen → Destino")
plt.tight_layout()
plt.savefig("charts/heatmap_rutas.png")
plt.close()

# ------------------------------------------
# CREAR POWERPOINT
# ------------------------------------------
ppt = Presentation()
title_slide = ppt.slides.add_slide(ppt.slide_layouts[0])
title_slide.shapes.title.text = "Reporte Under-Promise - Amazon México"
title_slide.placeholders[1].text = "Generado automáticamente en Python"

# 1. Métricas Generales
slide = ppt.slides.add_slide(ppt.slide_layouts[5])
shapes = slide.shapes
shapes.title.text = "Métricas Generales"
for i, row in enumerate(df_summary.values):
    text = f"{row[0]}: {row[1]}"
    txBox = shapes.add_textbox(Inches(0.5), Inches(1 + i*0.4), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    tf.text = text

# 2. Agregar gráficos
charts = ["destino", "origen", "sellers", "heatmap_rutas"]
titles = ["Destino", "Origen", "Top Sellers", "Rutas"]

for name, title in zip(charts, titles):
    slide = ppt.slides.add_slide(ppt.slide_layouts[5])
    slide.shapes.title.text = f"{title}"
    path = f"charts/{name}.png"
    slide.shapes.add_picture(path, Inches(1), Inches(1.3), width=Inches(8))

# Guardar presentación
ppt.save("underpromise_report.pptx")
print("✅ PowerPoint generado: underpromise_report.pptx")
